from pptx import Presentation
from gtts import gTTS
from docx import Document
from PyPDF2 import PdfReader
import os

def extract_text_from_pptx(file_path):
    try:
        prs = Presentation(file_path)
        slide_texts = []
        for slide in prs.slides:
            text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + " "
            slide_texts.append(text.strip())
        return slide_texts
    except FileNotFoundError:
        print(f"Error: File '{file_path}' not found.")
        return []
    except Exception as e:
        print(f"An error occurred while extracting text: {e}")
        return []
def extract_text_from_docx(file):
    doc = Document(file)
    return ["\n".join(para.text for para in doc.paragraphs if para.text.strip())]

def extract_text_from_pdf(file):
    reader = PdfReader(file)
    return [page.extract_text() for page in reader.pages if page.extract_text()]

def text_to_speech(text, output_path, lang='fr',tld='com'):
    try:
        tts = gTTS(text=text, lang=lang,tld=tld)
        tts.save(output_path)
    except Exception as e:
        print(f"An error occurred while generating audio: {e}")

